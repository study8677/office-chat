import os
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()

api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    raise RuntimeError("GOOGLE_API_KEY environment variable not set")

# Configure the Gemini client
genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-pro")

app = FastAPI(title="Office Chat")

class ChatRequest(BaseModel):
    message: str

class ChatResponse(BaseModel):
    response: str

@app.post("/chat", response_model=ChatResponse)
async def chat(req: ChatRequest):
    try:
        result = model.generate_content(req.message)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    return ChatResponse(response=result.text)


class DocumentRequest(BaseModel):
    kind: str
    prompt: str


class DocumentResponse(BaseModel):
    file: str


@app.post("/document", response_model=DocumentResponse)
async def create_document(req: DocumentRequest):
    if req.kind not in {"word", "excel", "ppt"}:
        raise HTTPException(status_code=400, detail="kind must be one of 'word', 'excel', 'ppt'")
    try:
        result = model.generate_content(req.prompt)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    text = result.text.strip()

    if req.kind == "word":
        doc = Document()
        doc.add_paragraph(text)
        filename = "output.docx"
        doc.save(filename)
    elif req.kind == "excel":
        wb = Workbook()
        ws = wb.active
        for i, line in enumerate(text.splitlines(), start=1):
            ws.cell(row=i, column=1, value=line)
        filename = "output.xlsx"
        wb.save(filename)
    else:  # ppt
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        body = slide.placeholders[1]
        tf = body.text_frame
        for line in text.splitlines():
            p = tf.add_paragraph()
            p.text = line
        filename = "output.pptx"
        prs.save(filename)

    return DocumentResponse(file=filename)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=8000)
